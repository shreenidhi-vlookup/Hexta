"""Text extraction from common document formats.

Supports PDF (pdfplumber), plain text, DOCX, HTML, Markdown, CSV,
XLSX (openpyxl), PPTX (python-pptx), ODT (zip + XML), RTF (striprtf),
and EPUB (zip + HTML). PDF text extraction is always available even
without OCR; OCR (Tesseract) is applied separately in ocr.py for
scanned PDFs.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

try:
    import pdfplumber

    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False
    logger.warning("pdfplumber not installed; PDF support disabled in batch mode")

try:
    import docx

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    logger.warning("python-docx not installed; DOCX support disabled")

try:
    import openpyxl

    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    logger.warning("openpyxl not installed; XLSX support disabled")

try:
    import pptx

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed; PPTX support disabled")

try:
    from striprtf.striprtf import rtf_to_text

    HAS_RTF = True
except ImportError:
    HAS_RTF = False
    logger.warning("striprtf not installed; RTF support disabled")


@dataclass
class ExtractedText:
    text: str
    pages: list[str]  # one entry per page (empty for plain text formats)
    source_format: str


def extract_text(file_path: str | Path) -> ExtractedText:
    """Extract text from a file, returning text + per-page breakdown."""
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(path)
    if ext in (".txt", ".md"):
        return _extract_text_file(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext in (".html", ".htm"):
        return _extract_html(path)
    if ext == ".csv":
        return _extract_csv(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".odt":
        return _extract_odt(path)
    if ext == ".rtf":
        return _extract_rtf(path)
    if ext == ".epub":
        return _extract_epub(path)

    raise ValueError(f"Unsupported file extension: {ext}")


def _extract_pdf(path: Path) -> ExtractedText:
    if not HAS_PDFPLUMBER:
        raise RuntimeError("pdfplumber is required for PDF extraction")
    pages = []
    full_text = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            pages.append(text)
            full_text.append(text)
    return ExtractedText(
        text="\n".join(full_text),
        pages=pages,
        source_format="pdf",
    )


def _extract_text_file(path: Path) -> ExtractedText:
    text = path.read_text(encoding="utf-8", errors="replace")
    return ExtractedText(text=text, pages=[text], source_format=path.suffix[1:])


def _extract_docx(path: Path) -> ExtractedText:
    if not HAS_DOCX:
        raise RuntimeError("python-docx is required for DOCX extraction")
    doc = docx.Document(str(path))
    pages = [para.text for para in doc.paragraphs]
    return ExtractedText(
        text="\n".join(p for p in pages if p),
        pages=pages,
        source_format="docx",
    )


def _extract_html(path: Path) -> ExtractedText:
    from html.parser import HTMLParser

    class TextExtractor(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.parts: list[str] = []

        def handle_data(self, data: str) -> None:
            self.parts.append(data)

    content = path.read_text(encoding="utf-8", errors="replace")
    parser = TextExtractor()
    parser.feed(content)
    text = "\n".join(parser.parts)
    return ExtractedText(text=text, pages=[text], source_format="html")


def _extract_csv(path: Path) -> ExtractedText:
    import csv

    rows: list[list[str]] = []
    with open(path, newline="", encoding="utf-8", errors="replace") as fh:
        reader = csv.reader(fh)
        for row in reader:
            rows.append([cell.strip() for cell in row])

    if not rows:
        return ExtractedText(text="", pages=[""], source_format="csv")

    text = _render_markdown_table(rows)
    return ExtractedText(text=text, pages=[text], source_format="csv")


def _extract_xlsx(path: Path) -> ExtractedText:
    if not HAS_OPENPYXL:
        raise RuntimeError("openpyxl is required for XLSX extraction")
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)

    blocks: list[str] = []
    for ws in wb.worksheets:
        rows = [
            ["" if cell is None else str(cell).strip() for cell in row]
            for row in ws.iter_rows(values_only=True)
            if any(cell is not None for cell in row)
        ]
        if not rows:
            continue
        blocks.append(f"Sheet: {ws.title}\n\n{_render_markdown_table(rows)}")

    wb.close()
    text = "\n\n".join(blocks)
    return ExtractedText(text=text, pages=[text], source_format="xlsx")


def _extract_pptx(path: Path) -> ExtractedText:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx is required for PPTX extraction")
    prs = pptx.Presentation(str(path))
    slides: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    para.text for para in shape.text_frame.paragraphs if para.text
                )
                if text:
                    parts.append(text)
            elif shape.shape_type == 19 and hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                parts.append(_render_markdown_table(rows))
        if parts:
            slides.append(f"--- Slide {i} ---\n" + "\n\n".join(parts))

    text = "\n\n".join(slides)
    return ExtractedText(text=text, pages=slides, source_format="pptx")


def _extract_odt(path: Path) -> ExtractedText:
    import xml.etree.ElementTree as ET
    import zipfile

    ns = {
        "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
        "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
    }
    with zipfile.ZipFile(str(path)) as zf:
        with zf.open("content.xml") as fh:
            tree = ET.parse(fh)
    root = tree.getroot()

    parts: list[str] = []
    for p in root.iter(f"{{{ns['text']}}}p"):
        text = "".join(p.itertext()).strip()
        if text:
            parts.append(text)

    for table in root.iter(f"{{{ns['table']}}}table"):
        rows: list[list[str]] = []
        for row in table.iter(f"{{{ns['table']}}}table-row"):
            cells = [
                "".join(cell.itertext()).strip()
                for cell in row.iter(f"{{{ns['table']}}}table-cell")
            ]
            if any(cells):
                rows.append(cells)
        if rows:
            parts.append(_render_markdown_table(rows))

    text = "\n\n".join(parts)
    return ExtractedText(text=text, pages=[text], source_format="odt")


def _extract_rtf(path: Path) -> ExtractedText:
    if not HAS_RTF:
        raise RuntimeError("striprtf is required for RTF extraction")
    content = path.read_text(encoding="utf-8", errors="replace")
    text = rtf_to_text(content)
    return ExtractedText(text=text, pages=[text], source_format="rtf")


def _extract_epub(path: Path) -> ExtractedText:
    import zipfile
    from html.parser import HTMLParser

    class TextExtractor(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.parts: list[str] = []

        def handle_data(self, data: str) -> None:
            self.parts.append(data)

    chapters: list[str] = []
    with zipfile.ZipFile(str(path)) as zf:
        for name in zf.namelist():
            if not name.lower().endswith((".xhtml", ".html", ".htm")):
                continue
            content = zf.read(name).decode("utf-8", errors="replace")
            parser = TextExtractor()
            parser.feed(content)
            chapter = "\n".join(parser.parts).strip()
            if chapter:
                chapters.append(chapter)

    text = "\n\n".join(chapters)
    return ExtractedText(text=text, pages=chapters, source_format="epub")


def _render_markdown_table(rows: list[list[str]]) -> str:
    """Render rows as a markdown table (pipe-separated with a separator row)."""
    width = max((len(r) for r in rows), default=1) or 1
    padded = [(r + [""] * width)[:width] for r in rows]
    sep = "| " + " | ".join(["---"] * width) + " |"
    lines = ["| " + " | ".join(row) + " |" for row in padded]
    lines.insert(1, sep)
    return "\n".join(lines)
