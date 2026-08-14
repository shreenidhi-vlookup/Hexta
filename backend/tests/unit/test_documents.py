"""Unit tests for document ingestion modules."""

from __future__ import annotations

from app.documents import categories
from app.documents.entity_extraction import extract_entities
from app.documents.validation import validate_upload
from app.documents.chunking.structural_chunker import StructuralChunker
from app.documents.text_extraction import ExtractedText, extract_text
from app.documents.metadata_extraction import extract_metadata


class TestValidation:
    def test_valid_txt_file(self):
        result = validate_upload("doc.txt", 1024)
        assert result.valid is True
        assert result.error is None

    def test_valid_pdf_file(self):
        result = validate_upload("doc.pdf", 1024 * 1024)
        assert result.valid is True

    def test_invalid_extension(self):
        result = validate_upload("doc.exe", 100)
        assert result.valid is False
        assert "not allowed" in result.error

    def test_oversized_file(self):
        result = validate_upload("doc.txt", 25 * 1024 * 1024)
        assert result.valid is False
        assert "too big" in result.error
        assert "20 MB" in result.error
        assert "smaller than" in result.error

    def test_case_insensitive_extension(self):
        result = validate_upload("DOC.PDF", 1024)
        assert result.valid is True

    def test_rejects_path_traversal_unix_style(self):
        result = validate_upload("../../../etc/passwd.txt", 1024)
        assert result.valid is False
        assert "path separators" in result.error

    def test_rejects_path_traversal_windows_style(self):
        result = validate_upload("..\\..\\windows\\system32\\evil.txt", 1024)
        assert result.valid is False
        assert "path separators" in result.error

    def test_rejects_embedded_slash(self):
        result = validate_upload("storage/processed/overwrite.txt", 1024)
        assert result.valid is False
        assert "path separators" in result.error

    def test_rejects_null_byte(self):
        result = validate_upload("doc.txt\x00.pdf", 1024)
        assert result.valid is False

    def test_rejects_bare_dotdot(self):
        result = validate_upload("..", 1024)
        assert result.valid is False

    def test_plain_filename_with_dots_still_valid(self):
        # A dotted-but-otherwise-normal filename must not be caught by the
        # traversal check.
        result = validate_upload("v1.2.final_report.pdf", 1024)
        assert result.valid is True


class TestStructuralChunker:
    def test_chunk_plain_text(self):
        extractor = ExtractedText(
            text="This is a test. " * 100,  # ~230 words
            pages=[],
            source_format="txt",
        )
        chunker = StructuralChunker(max_tokens=50)
        chunks = list(chunker.chunk(extractor))
        assert len(chunks) > 1
        assert all(c.chunk_type == "paragraph" for c in chunks)

    def test_chunk_preserves_tables(self):
        text = (
            "Some intro paragraph here.\n\n"
            "col1 col2 col3\n"
            "val1 val2 val3\n"
            "val4 val5 val6\n\n"
            "Another paragraph after the table."
        )
        extractor = ExtractedText(text=text, pages=[], source_format="txt")
        chunker = StructuralChunker(max_tokens=500)
        chunks = list(chunker.chunk(extractor))

        table_chunks = [c for c in chunks if c.chunk_type == "table"]
        assert len(table_chunks) >= 1
        assert "col1" in table_chunks[0].content

    def test_chunk_empty_text(self):
        extractor = ExtractedText(text="", pages=[], source_format="txt")
        chunker = StructuralChunker()
        chunks = list(chunker.chunk(extractor))
        assert len(chunks) == 0


class TestMetadataExtraction:
    def test_extract_title_from_content(self):
        text = "Housing Finance Guidelines\n\nSection 1: Credit scores..."
        extractor = ExtractedText(text=text, pages=[], source_format="txt")
        meta = extract_metadata(extractor, "doc.txt")
        assert meta.title is not None
        assert meta.doc_type == "policy"

    def test_extract_doc_type_from_content(self):
        text = "Underwriting guidelines for conventional loans..."
        extractor = ExtractedText(text=text, pages=[], source_format="txt")
        meta = extract_metadata(extractor, "underwriting.txt")
        assert meta.doc_type == "underwriting"


class TestEntityExtraction:
    def test_extracts_lenders(self):
        entities = extract_entities(
            "The FHA loan guidelines and VA loan program details."
        )
        assert "federal housing administration" in entities.lenders
        assert "veterans affairs" in entities.lenders

    def test_extracts_products(self):
        entities = extract_entities("An adjustable rate mortgage ARM or HELOC option.")
        assert "adjustable rate mortgage" in entities.products
        assert "helic home equity line of credit" in entities.products

    def test_extracts_documents(self):
        entities = extract_entities("Submit two pay stubs and a W-2 with your application.")
        assert "pay stub" in entities.documents
        assert "w-2" in entities.documents

    def test_extracts_property_types(self):
        entities = extract_entities("Financing a primary residence or an investment property.")
        assert "primary residence" in entities.property_types
        assert "investment property" in entities.property_types

    def test_empty_text_no_entities(self):
        entities = extract_entities("")
        assert entities.lenders == []
        assert entities.products == []
        assert entities.documents == []
        assert entities.property_types == []
        assert entities.acronyms == []


class TestTextExtractionFormats:
    def test_extract_csv(self, tmp_path):
        p = tmp_path / "data.csv"
        p.write_text("lender,limit\ngeneral,5000\nunderwriting,25000\n", encoding="utf-8")
        result = extract_text(p)
        assert "lender" in result.text
        assert "5000" in result.text
        assert "---" in result.text

    def test_extract_xlsx(self, tmp_path):
        import openpyxl

        p = tmp_path / "limits.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["department", "approval_limit"])
        ws.append(["general", "5000"])
        wb.save(str(p))
        result = extract_text(p)
        assert "department" in result.text
        assert "5000" in result.text

    def test_extract_pptx(self, tmp_path):
        from pptx import Presentation

        p = tmp_path / "slides.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Approval Limits"
        prs.save(str(p))
        result = extract_text(p)
        assert "Approval Limits" in result.text

    def test_extract_rtf(self, tmp_path):
        p = tmp_path / "notes.rtf"
        p.write_bytes(b"{\\rtf1\\ansi Approval limit is \\b GBP 5000\\b0 .}")
        result = extract_text(p)
        assert "Approval limit" in result.text

    def test_extract_odt(self, tmp_path):
        import zipfile

        p = tmp_path / "doc.odt"
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<office:document-content '
            'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
            'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" '
            'xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0">'
            '<office:body><office:text>'
            "<text:p>Claims approval limit</text:p>"
            "<text:p>Escalate above limit</text:p>"
            "</office:text></office:body></office:document-content>"
        )
        with zipfile.ZipFile(str(p), "w") as zf:
            zf.writestr("content.xml", content_xml)
        result = extract_text(p)
        assert "Claims approval limit" in result.text

    def test_extract_epub(self, tmp_path):
        import zipfile

        p = tmp_path / "book.epub"
        chapter = '<html><body><p>Life insurance basics</p></body></html>'
        with zipfile.ZipFile(str(p), "w") as zf:
            zf.writestr("EPUB/chapter1.xhtml", chapter)
        result = extract_text(p)
        assert "Life insurance basics" in result.text

    def test_extract_htm(self, tmp_path):
        p = tmp_path / "page.htm"
        p.write_text("<html><body><h1>Guide</h1><p>Approval workflow.</p></body></html>", encoding="utf-8")
        result = extract_text(p)
        assert "Approval workflow" in result.text

    def test_unsupported_extension_raises(self, tmp_path):
        p = tmp_path / "archive.zip"
        p.write_bytes(b"not a document")
        try:
            extract_text(p)
        except ValueError:
            return
        raise AssertionError("expected ValueError for unsupported extension")


class TestAllowedExtensions:
    def test_new_formats_accepted(self):
        for name in [
            "a.csv", "a.xlsx", "a.pptx", "a.odt", "a.rtf", "a.epub", "a.htm",
        ]:
            assert validate_upload(name, 1024).valid is True

    def test_unknown_extension_rejected(self):
        result = validate_upload("a.zip", 1024)
        assert result.valid is False
        assert "not allowed" in result.error


class TestAutoIngest:
    def test_trigger_spawns_ingest_batch(self, monkeypatch):
        import subprocess

        from app.documents import auto_ingest

        spawned = {}

        def fake_popen(cmd, **kwargs):
            spawned["cmd"] = cmd
            spawned["kwargs"] = kwargs
            return object()

        monkeypatch.setattr(subprocess, "Popen", fake_popen)

        ok = auto_ingest.trigger_ingestion("storage/pending")

        assert ok is True
        assert "-m" in spawned["cmd"]
        assert "app.documents.ingest_batch" in spawned["cmd"]
        assert "--queue-dir" in spawned["cmd"]
        assert "storage/pending" in spawned["cmd"]
        # PYTHONPATH must include the backend root so `-m` resolves
        assert "PYTHONPATH" in spawned["kwargs"]["env"]

    def test_trigger_handles_spawn_failure(self, monkeypatch):
        import subprocess

        from app.documents import auto_ingest

        def boom(cmd, **kwargs):
            raise RuntimeError("spawn failed")

        monkeypatch.setattr(subprocess, "Popen", boom)

        assert auto_ingest.trigger_ingestion("storage/pending") is False

    def test_backend_root_points_to_backend(self):
        from app.documents.auto_ingest import _backend_root

        root = _backend_root()
        assert root.name == "backend"
        assert (root / "app").is_dir()
        assert (root / "app" / "documents" / "ingest_batch.py").exists()


class TestMetadataCategoryOverride:
    """An explicit choice at upload beats content detection.

    Detection reads keywords out of the first 2000 characters, which
    misfiles anything that merely mentions a category -- the glossary in
    the knowledge base was typed "underwriting" because one of its
    definitions uses the word.
    """

    def _extracted(self, text: str) -> ExtractedText:
        return ExtractedText(text=text, pages=[text], source_format="txt")

    GLOSSARY = (
        "General Banking and Mortgage Glossary\n\n"
        "Underwriting: The process by which a lender evaluates risk.\n"
    )

    def test_detection_still_applies_without_an_override(self):
        meta = extract_metadata(self._extracted(self.GLOSSARY), "g.txt")
        assert meta.doc_type == "underwriting"  # the misfiling, unchanged

    def test_explicit_doc_type_wins(self):
        meta = extract_metadata(
            self._extracted(self.GLOSSARY), "g.txt", doc_type="glossary",
        )
        assert meta.doc_type == "glossary"

    def test_explicit_department_is_used(self):
        meta = extract_metadata(
            self._extracted(self.GLOSSARY), "g.txt", department="compliance",
        )
        assert meta.department == "compliance"

    def test_department_defaults_when_absent(self):
        meta = extract_metadata(self._extracted(self.GLOSSARY), "g.txt")
        assert meta.department == categories.DEFAULT_DEPARTMENT

    def test_auto_sentinel_falls_back_to_detection(self):
        meta = extract_metadata(
            self._extracted(self.GLOSSARY),
            "g.txt",
            doc_type=categories.AUTO_DOC_TYPE,
        )
        assert meta.doc_type == "underwriting"

    def test_invalid_override_falls_back_rather_than_raising(self):
        """The endpoint rejects bad values; a hand-edited sidecar must not
        be able to crash a batch run."""
        meta = extract_metadata(
            self._extracted(self.GLOSSARY),
            "g.txt",
            doc_type="invoice",
            department="marketing",
        )
        assert meta.doc_type == "underwriting"
        assert meta.department == categories.DEFAULT_DEPARTMENT
